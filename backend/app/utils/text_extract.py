import io
import re
from typing import Optional

import docx
from pypdf import PdfReader
from pptx import Presentation


def _clean_text(raw: str) -> str:
    """
    Normalize whitespace and strip control characters to make downstream chunking more stable.
    """
    text = raw.replace("\r\n", "\n")
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_text_from_bytes(data: bytes, filename: str) -> str:
    lower = filename.lower()

    if lower.endswith(".txt") or lower.endswith(".md"):
        return _clean_text(data.decode("utf-8", errors="ignore"))

    if lower.endswith(".pdf"):
        reader = PdfReader(io.BytesIO(data))
        parts = []
        for page in reader.pages:
            parts.append(page.extract_text() or "")
        return _clean_text("\n\n".join(parts))

    if lower.endswith(".docx"):
        d = docx.Document(io.BytesIO(data))
        return _clean_text("\n".join([p.text for p in d.paragraphs]))

    if lower.endswith(".pptx"):
        prs = Presentation(io.BytesIO(data))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return _clean_text("\n".join(texts))

    # Fallback: treat as text
    return _clean_text(data.decode("utf-8", errors="ignore"))
